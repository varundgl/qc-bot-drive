# src/preprocessing/file_processor.py
from pptx import Presentation
import nbformat
import re
import os
import logging

logger = logging.getLogger(__name__)

class FileProcessor:
    @staticmethod
    def process_slide_file(file_path: str) -> str:
        try:
            prs = Presentation(file_path)
            content = []
            for i, slide in enumerate(prs.slides):
                content.append(f"=== Slide {i+1} ===")
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        cleaned_text = re.sub(r'\s+', ' ', shape.text.strip())
                        content.append(cleaned_text)
                content.append("")
            return "\n".join(content)
        except Exception as e:
            logger.error(f"Error processing presentation: {str(e)}")
            return ""

    @staticmethod
    def process_notebook_file(file_path: str) -> str:
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                nb = nbformat.read(f, as_version=4)

            content = []
            for cell in nb.cells:
                if cell.cell_type == 'code':
                    content.append("## CODE CELL ##")
                    content.append(cell.source.strip())
                    content.append("----")
                elif cell.cell_type == 'markdown':
                    content.append("## MARKDOWN CELL ##")
                    cleaned_text = cell.source.strip()
                    cleaned_text = re.sub(r'#+\s*', '', cleaned_text)
                    cleaned_text = re.sub(r'\*{1,2}(.*?)\*{1,2}', r'\1', cleaned_text)
                    cleaned_text = re.sub(r'\[(.*?)\]\(.*?\)', r'\1', cleaned_text)
                    content.append(cleaned_text)
                    content.append("----")

            return "\n".join(content)
        except Exception as e:
            logger.error(f"Error processing notebook: {str(e)}")
            return ""